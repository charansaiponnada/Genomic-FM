"""
Build the Mnemosyne research deck (python-pptx). Genomics-themed palette tied to
the architecture figure. Honest results (verified vs. pending), 25 real references.
Output: deck/Mnemosyne_Research_Deck.pptx
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = os.path.dirname(os.path.abspath(__file__))
ASSET = lambda n: os.path.join(HERE, "assets", n)

# ---------- palette ----------
NAVY   = RGBColor(0x0F, 0x22, 0x33)
NAVY2  = RGBColor(0x16, 0x2E, 0x40)
TEAL   = RGBColor(0x2E, 0x8B, 0x7A)
BLUE   = RGBColor(0x3A, 0x6E, 0xA5)
BLUEL  = RGBColor(0x5B, 0x8D, 0xEF)
AMBER  = RGBColor(0xE0, 0x8A, 0x2B)
INK    = RGBColor(0x1B, 0x1B, 0x24)
MUTED  = RGBColor(0x6A, 0x6A, 0x78)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)
PALE   = RGBColor(0xF2, 0xF6, 0xF8)
PALEB  = RGBColor(0xE9, 0xF0, 0xF7)

HEAD = "Cambria"
BODY = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True, space=1.0, space_after=4):
    """runs: list of paragraphs; each paragraph is list of (txt,size,bold,color,font)."""
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(1)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = space
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        for (txt, size, bold, color, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = font
    return tb


def card(s, l, t, w, h, fill=PALE, line=None, radius=0.08):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t),
                            Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(1)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = radius
    except Exception:
        pass
    return sh


def chip(s, l, t, d, label, fill, tcolor=WHITE):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = fill; c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label
    r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = tcolor; r.font.name = HEAD
    return c


def title(s, txt, color=INK, sub=None, subcolor=TEAL):
    text(s, 0.7, 0.5, 12.0, 1.0, [[(txt, 34, True, color, HEAD)]])
    if sub:
        text(s, 0.72, 1.28, 12.0, 0.5, [[(sub, 15, False, subcolor, BODY)]])


def dot(s, l, t, color, d=0.16):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    c.shadow.inherit = False
    return c


def page_num(s, n, dark=False):
    text(s, 12.4, 7.05, 0.8, 0.3, [[(str(n), 10, False, ICE if dark else MUTED, BODY)]],
         align=PP_ALIGN.RIGHT)


# ============================== 1. TITLE ==============================
s = slide(NAVY)
text(s, 0.9, 2.2, 11.5, 1.4, [[("Mnemosyne", 60, True, WHITE, HEAD)]])
text(s, 0.95, 3.5, 11.5, 0.9,
     [[("Linear-Time Genomic Sequence Modeling with a Unified Associative Memory",
        22, False, ICE, BODY)]])
text(s, 0.95, 4.25, 11.5, 0.5,
     [[("Persistent motifs  +  sequence registers   ·   reverse-complement equivariant",
        14, False, TEAL, BODY)]])
text(s, 0.95, 5.6, 11.5, 0.5, [[("Charan Sai Ponnada", 17, True, WHITE, BODY)]])
text(s, 0.95, 6.05, 11.5, 0.5, [[("Research overview  ·  July 2026", 12, False, MUTED, BODY)]])
for i, c in enumerate((TEAL, BLUE, AMBER)):
    dot(s, 0.95 + i * 0.34, 5.05, c, 0.2)

# ============================== 2. PROBLEM ==============================
s = slide(WHITE)
title(s, "Problem statement", sub="DNA function spans enormous distances — a model must satisfy three demands at once")
reqs = [("O(L)", "Linear-time scaling", "Reach 10^5–10^6 bp enhancer–promoter contexts at all.", TEAL),
        ("A↔", "Content-addressable recall", "Retrieve sparse motifs & distal elements by content, not position.", BLUE),
        ("⟲", "Reverse-complement symmetry", "DNA is double-stranded: forward and RC strands are the same molecule.", AMBER)]
x = 0.7
for icon, head, body, col in reqs:
    card(s, x, 1.95, 3.85, 2.5, PALE)
    chip(s, x + 0.3, 2.2, 0.7, icon, col)
    text(s, x + 0.25, 3.15, 3.35, 0.6, [[(head, 15, True, INK, BODY)]])
    text(s, x + 0.25, 3.62, 3.4, 1.0, [[(body, 12, False, MUTED, BODY)]], space=1.02)
    x += 4.05
card(s, 0.7, 4.75, 11.9, 1.9, PALEB)
text(s, 1.0, 4.95, 11.4, 1.6,
     [[("No mainstream architecture does all three.", 16, True, INK, BODY)],
      [("• Transformers (DNABERT-2): strong recall, but O(L²) — context is capped.",
        13.5, False, INK, BODY)],
      [("• Selective SSMs (Mamba, HyenaDNA): O(L), but a fixed-width state means recall collapses with distance.",
        13.5, False, INK, BODY)],
      [("• Caduceus: adds the missing RC symmetry to Mamba — but inherits the fixed-state recall bottleneck.",
        13.5, False, INK, BODY)]], space=1.05, space_after=3)
page_num(s, 2)

# ============================== 3. OBJECTIVES ==============================
s = slide(WHITE)
title(s, "Objectives", sub="Turn the three demands into one linear-cost mechanism — and test it honestly")
objs = [
    ("Add a linear-cost associative memory that every token reads by one Hopfield step.", TEAL),
    ("Unify two slot types: a learned universal-motif dictionary (persistent) and per-sequence long-range recall (registers).", BLUE),
    ("Guarantee exact reverse-complement equivariance by construction.", AMBER),
    ("Prove linear time (semiseparable + low-rank) and verify the memory is actually active.", TEAL),
    ("Test the falsifiable claim — memory beats a matched SSM at long range — with no fabricated numbers.", BLUE),
]
y = 1.95
for i, (o, col) in enumerate(objs, 1):
    chip(s, 0.8, y, 0.62, str(i), col)
    card(s, 1.65, y - 0.03, 10.9, 0.78, PALE)
    text(s, 1.95, y + 0.02, 10.4, 0.7, [[(o, 14.5, False, INK, BODY)]],
         anchor=MSO_ANCHOR.MIDDLE, space=1.0)
    y += 0.96
page_num(s, 3)

# ============================== 4. LIT REVIEW I ==============================
s = slide(WHITE)
title(s, "Related work I — sequence models & associative memory", subcolor=BLUE)
groups = [
    (TEAL, "State-space models", "S4 (Gu '22) → Mamba (Gu & Dao '23) → Mamba-2 / SSD (Dao & Gu '24): linear-time, but a known recall bottleneck."),
    (TEAL, "Long convolutions", "Hyena (Poli '23), HyenaDNA (Nguyen '23): million-token contexts via implicit convolutions."),
    (BLUE, "Hybrids re-add attention", "Jamba (Lieber '24), Griffin (De '24): interleave attention to recover recall — at higher cost."),
    (BLUE, "The recall gap", "Zoology / MQAR (Arora '23), Based (Arora '24): efficient models under-recall; MQAR isolates the gap."),
    (AMBER, "Associative memory", "Modern Hopfield (Ramsauer '21), Energy Transformer (Hoover '23), optimal capacity (Wu '24), Perceiver bottleneck (Jaegle '21)."),
    (AMBER, "Pre-training", "Attention Is All You Need (Vaswani '17); Masked Autoencoders (He '22) — our pre-training objective."),
]
x, y = 0.7, 1.85
for i, (col, head, body) in enumerate(groups):
    cx = x + (i % 2) * 6.05
    cy = y + (i // 2) * 1.62
    card(s, cx, cy, 5.85, 1.45, PALE)
    dot(s, cx + 0.28, cy + 0.32, col, 0.2)
    text(s, cx + 0.62, cy + 0.18, 5.1, 0.5, [[(head, 14.5, True, INK, BODY)]])
    text(s, cx + 0.3, cy + 0.66, 5.35, 0.75, [[(body, 11.5, False, MUTED, BODY)]], space=1.02)
page_num(s, 4)

# ============================== 5. LIT REVIEW II ==============================
s = slide(WHITE)
title(s, "Related work II — genomic foundation models & benchmarks", subcolor=BLUE)
g2 = [
    (BLUE, "Transformer DNA LMs", "DNABERT (Ji '21), DNABERT-2 (Zhou '24), Nucleotide Transformer (Dalla-Torre '25), GENA-LM (Fishman '25), GROVER ('24)."),
    (TEAL, "Efficient / long-range", "HyenaDNA (Nguyen '23); Caduceus (Schiff '24) — RC-equivariant Mamba, the closest prior design."),
    (AMBER, "Frontier generative", "Evo (Nguyen '24, 7B); Evo 2 (Brixi '25, 40B, 9.3T nt, 1 Mb ctx) — the scale frontier."),
    (AMBER, "Regulatory / functional", "Enformer (Avsec '21), Borzoi (Linder '25), AlphaGenome (Avsec '25) — variant-effect SOTA."),
    (TEAL, "Benchmarks", "GUE, Genomic Benchmarks, BEND (Marin '24) — where we will evaluate, fairly."),
    (BLUE, "The gap we target", "None combines O(L) + content-addressable recall + RC symmetry. Mnemosyne aims exactly there."),
]
x, y = 0.7, 1.85
for i, (col, head, body) in enumerate(g2):
    cx = x + (i % 2) * 6.05
    cy = y + (i // 2) * 1.62
    fill = PALEB if head == "The gap we target" else PALE
    card(s, cx, cy, 5.85, 1.45, fill)
    dot(s, cx + 0.28, cy + 0.32, col, 0.2)
    text(s, cx + 0.62, cy + 0.18, 5.1, 0.5, [[(head, 14.5, True, INK, BODY)]])
    text(s, cx + 0.3, cy + 0.66, 5.35, 0.75, [[(body, 11.5, False, MUTED, BODY)]], space=1.02)
page_num(s, 5)

# ============================== 6. METHOD 1: ARCHITECTURE ==============================
s = slide(WHITE)
title(s, "Methodology (1/3) — Architecture", sub="A selective SSM in parallel with a unified associative memory, fused by a learned gate")
_ih = 4.95
_iw = _ih * 2934.0 / 1993.0
s.shapes.add_picture(ASSET("architecture.png"), Inches((13.333 - _iw) / 2), Inches(1.7),
                     height=Inches(_ih))
text(s, 0.7, 6.85, 11.9, 0.5,
     [[("Left: RC-equivariant masked-autoencoder stack.   Right: one Mnemosyne block — token-mixing = semiseparable (SSM) + low-rank (memory), hence O(L).",
        11, False, MUTED, BODY)]], align=PP_ALIGN.CENTER)
page_num(s, 6)

# ============================== 7. METHOD 2: MEMORY + THEORY ==============================
s = slide(WHITE)
title(s, "Methodology (2/3) — Unified associative memory", subcolor=BLUE)
# left column: description
text(s, 0.7, 1.9, 6.0, 4.6,
     [[("Every token reads S = P_p + P_r slots by one Hopfield step:", 14, True, INK, BODY)],
      [("", 6, False, INK, BODY)],
      [("Persistent slots  M", 13.5, True, TEAL, BODY)],
      [("Learned parameters — a differentiable dictionary of universal motifs, shared across all sequences.",
        12.5, False, MUTED, BODY)],
      [("", 5, False, INK, BODY)],
      [("Register slots", 13.5, True, BLUE, BODY)],
      [("Gathered from the current sequence (Perceiver bottleneck) → per-sequence, distance-independent recall.",
        12.5, False, MUTED, BODY)],
      [("", 5, False, INK, BODY)],
      [("Learned per-channel gate", 13.5, True, AMBER, BODY)],
      [("σ(·)⊙ — starts closed (negative-bias init), opens as memory earns its place. Cannot be silenced by one logit (v1's failure).",
        12.5, False, MUTED, BODY)]], space=1.05)
# right column: equations + theorem
card(s, 7.0, 1.9, 5.6, 1.55, PALEB)
text(s, 7.3, 2.05, 5.1, 1.3,
     [[("Read:   rᵢ = V · softmax( β · qᵢ Kᵀ )", 14, True, INK, "Cambria")],
      [("Gather:  Reg = softmax( β Z Xᵀ ) · X", 14, True, INK, "Cambria")],
      [("β learnable → soft motif-family vs. sharp single-motif recall", 11, False, MUTED, BODY)]], space=1.15)
card(s, 7.0, 3.65, 5.6, 1.5, PALE, line=TEAL)
text(s, 7.3, 3.8, 5.1, 1.25,
     [[("Theorem (linear time)", 13.5, True, TEAL, BODY)],
      [("token-mixing = N-semiseparable (SSM) + rank ≤ P_p+P_r (memory)", 12, False, INK, BODY)],
      [("⇒ O( L · (N + P_p + P_r) · d ),  no L×L matrix formed.", 12, True, INK, BODY)]], space=1.08)
card(s, 7.0, 5.32, 5.6, 1.2, PALE, line=AMBER)
text(s, 7.3, 5.46, 5.1, 1.0,
     [[("Reverse-complement equivariance", 13.5, True, AMBER, BODY)],
      [("output = ½[ f(x) + S·f(Rx) ] — exact; tested to 0.0 error.", 12, False, INK, BODY)]], space=1.08)
page_num(s, 7)

# ============================== 8. METHOD 3: DATA & PROTOCOL ==============================
s = slide(WHITE)
title(s, "Methodology (3/3) — Data & experimental protocol", subcolor=BLUE)
# data collection flow
text(s, 0.7, 1.75, 11.5, 0.4, [[("Data collection", 15, True, INK, BODY)]])
flow = ["Multi-species genomes\n(UCSC: human, mouse,\nzebrafish, …)",
        "Sliding windows\n1024 bp, stride 512,\nfilter >5% N",
        "Single-nucleotide\ntokenize\n{A,C,G,T,N}",
        "MAE objective\nmask 30%,\npredict originals"]
x = 0.7
for i, f in enumerate(flow):
    card(s, x, 2.2, 2.65, 1.15, PALE, line=TEAL if i % 2 == 0 else BLUE)
    text(s, x + 0.15, 2.32, 2.4, 0.95, [[(f, 11.5, False, INK, BODY)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=1.0)
    if i < 3:
        text(s, x + 2.62, 2.5, 0.5, 0.6, [[("→", 22, True, MUTED, BODY)]])
    x += 3.05
# experiments
text(s, 0.7, 3.75, 11.5, 0.4, [[("Experiments (controlled — SSM baseline = same backbone, memory OFF)", 15, True, INK, BODY)]])
exps = [
    ("MQAR long-range", "Synthetic recall at controllable distance — the falsifiable core test.", TEAL),
    ("Controlled MAE", "Memory ON vs OFF; log gate + Hopfield energy to prove memory is active.", BLUE),
    ("GUE downstream", "Frozen linear probe on BOTH models (fair); report MCC. No probe-vs-finetune spin.", AMBER),
    ("Ablations", "Sweep P_r, P_p, β, and gate to show which component earns its keep.", TEAL),
]
x, y = 0.7, 4.25
for i, (h, b, col) in enumerate(exps):
    cx = x + (i % 2) * 6.05
    cy = y + (i // 2) * 1.15
    card(s, cx, cy, 5.85, 1.0, PALE)
    dot(s, cx + 0.25, cy + 0.22, col, 0.18)
    text(s, cx + 0.58, cy + 0.1, 5.1, 0.4, [[(h, 13, True, INK, BODY)]])
    text(s, cx + 0.3, cy + 0.5, 5.4, 0.45, [[(b, 11, False, MUTED, BODY)]], space=1.0)
text(s, 0.7, 6.65, 11.9, 0.4,
     [[("All paper numbers are auto-generated from run logs; any unfilled result renders as a red TODO.",
        11, True, TEAL, BODY)]])
page_num(s, 8)

# ============================== 9. RESULTS ==============================
s = slide(WHITE)
title(s, "Results — verified now, and the falsifiable prediction", subcolor=BLUE)
stats = [("7 / 7", "CPU correctness tests pass", TEAL),
         ("0.0", "RC-equivariance error (exact)", AMBER),
         ("gate 0.12 · E −27", "memory is ACTIVE (not inert — v1's failure)", BLUE),
         ("27M vs 16M", "memory vs matched SSM baseline (controlled)", TEAL)]
y = 1.95
for val, lab, col in stats:
    card(s, 0.7, y, 5.7, 1.05, PALE)
    text(s, 0.95, y + 0.06, 5.2, 0.6, [[(val, 24, True, col, HEAD)]])
    text(s, 0.98, y + 0.62, 5.2, 0.4, [[(lab, 12, False, MUTED, BODY)]])
    y += 1.18
s.shapes.add_picture(ASSET("expected_longrange.png"), Inches(6.7), Inches(1.95),
                     width=Inches(5.9))
text(s, 6.7, 6.35, 5.9, 0.9,
     [[("Prediction: SSM recall decays with distance; memory stays flat.", 12, True, INK, BODY)],
      [("Headline MQAR / GUE numbers pending free-GPU runs — we report only what we measure.",
        11, False, MUTED, BODY)]], space=1.05)
page_num(s, 9)

# ============================== 10. ASSESSMENT ==============================
s = slide(NAVY)
text(s, 0.7, 0.5, 12, 1.0, [[("Is it working? Groundbreaking? Why might it fail?", 30, True, WHITE, HEAD)]])
cols = [
    (TEAL, "Working?", ["Mechanism runs and is active — verified.",
                        "RC-equivariance proven (0.0 error).",
                        "The scientific claim is a ready hypothesis, not yet a measured result."]),
    (BLUE, "Groundbreaking?", ["Not by out-scaling — ~1000× smaller than Evo 2 / AlphaGenome.",
                        "Yes if the mechanism is rigorously shown to help and others adopt it — as Caduceus did with one idea.",
                        "Wrong bar: 'like AIAYN' = a simple mechanism that scales."]),
    (AMBER, "Why it might fail", ["MQAR may not separate from a tuned SSM.",
                        "Gate may collapse to 0 at scale.",
                        "Persistent slots may be dead weight; registers have a capacity ceiling.",
                        "May not transfer to genome scale."]),
]
x = 0.7
for col, head, items in cols:
    card(s, x, 1.7, 3.95, 5.0, NAVY2)
    chip(s, x + 0.3, 1.95, 0.5, "", col)
    text(s, x + 0.95, 2.02, 3.0, 0.5, [[(head, 17, True, WHITE, BODY)]])
    runs = []
    for it in items:
        runs.append([("•  ", 12, True, col, BODY), (it, 12, False, ICE, BODY)])
    text(s, x + 0.3, 2.75, 3.45, 3.8, runs, space=1.06, space_after=8)
    x += 4.05
page_num(s, 10, dark=True)

# ============================== 11. ROADMAP ==============================
s = slide(WHITE)
title(s, "Roadmap", sub="Cheapest decisive signal first, then scale — the real test of the mechanism")
steps = [
    ("1", "Run MQAR + ablations on a free GPU (Modal). The cheap, falsifiable signal.", TEAL),
    ("2", "Controlled MAE pre-training + fair GUE probes; confirm memory stays active.", BLUE),
    ("3", "Scale: human genome → multi-species; 27M → 350M. Show the memory benefit GROWS with scale.", AMBER),
    ("4", "Full GUE (28 tasks) + NT benchmarks vs current baselines; 3–5 seeds + confidence intervals.", BLUE),
    ("5", "Target a niche SOTA (splice / non-coding variant) where long-range recall is the right tool.", TEAL),
]
y = 1.95
for n, txt_, col in steps:
    chip(s, 0.8, y, 0.66, n, col)
    card(s, 1.7, y - 0.02, 10.85, 0.82, PALE)
    text(s, 2.0, y + 0.04, 10.3, 0.7, [[(txt_, 14, False, INK, BODY)]],
         anchor=MSO_ANCHOR.MIDDLE, space=1.0)
    y += 0.98
page_num(s, 11)

# ============================== 12. REFERENCES ==============================
s = slide(WHITE)
title(s, "References (selected, 25)", subcolor=BLUE)
refs = [
    "Vaswani et al. Attention Is All You Need. NeurIPS 2017.",
    "He et al. Masked Autoencoders Are Scalable Vision Learners. CVPR 2022.",
    "Gu et al. Efficiently Modeling Long Sequences with Structured State Spaces (S4). ICLR 2022.",
    "Gu & Dao. Mamba: Linear-Time Sequence Modeling with Selective State Spaces. 2023.",
    "Dao & Gu. Transformers are SSMs: Structured State Space Duality (Mamba-2). ICML 2024.",
    "Poli et al. Hyena Hierarchy: Larger Convolutional Language Models. ICML 2023.",
    "Lieber et al. Jamba: A Hybrid Transformer-Mamba Language Model. 2024.",
    "De et al. Griffin: Gated Linear Recurrences + Local Attention. 2024.",
    "Ramsauer et al. Hopfield Networks is All You Need. ICLR 2021.",
    "Hoover et al. Energy Transformer. NeurIPS 2023.",
    "Wu et al. Provably Optimal Memory Capacity for Modern Hopfield Models. NeurIPS 2024.",
    "Jaegle et al. Perceiver: General Perception with Iterative Attention. ICML 2021.",
    "Arora et al. Zoology: Measuring and Improving Recall in Efficient LMs. 2023.",
    "Arora et al. Based: Balancing the Recall-Throughput Tradeoff. ICML 2024.",
    "Ji et al. DNABERT. Bioinformatics 2021.",
    "Zhou et al. DNABERT-2: Efficient Multi-Species Genome FM. ICLR 2024.",
    "Dalla-Torre et al. The Nucleotide Transformer. Nature Methods 2025.",
    "Nguyen et al. HyenaDNA: Single-Nucleotide Long-Range Modeling. NeurIPS 2023.",
    "Schiff et al. Caduceus: Bi-Directional Equivariant DNA Modeling. ICML 2024.",
    "Nguyen et al. Sequence Modeling & Design at Genome Scale (Evo). Science 2024.",
    "Brixi et al. Genome Modeling & Design across Life with Evo 2. 2025.",
    "Avsec et al. Gene Expression Prediction from Sequence (Enformer). Nat. Methods 2021.",
    "Linder et al. Predicting RNA-seq Coverage at Base Resolution (Borzoi). Nat. Genet. 2025.",
    "Avsec et al. AlphaGenome: Regulatory Variant-Effect Prediction. Nature 2025.",
    "Marin et al. BEND: Benchmarking DNA LMs on Meaningful Tasks. ICLR 2024.",
]
col_w, x0, y0 = 6.0, 0.7, 1.8
for i, r in enumerate(refs):
    cx = x0 + (i // 13) * 6.15
    cy = y0 + (i % 13) * 0.395
    text(s, cx, cy, col_w, 0.4,
         [[(f"{i+1}. ", 9.5, True, BLUE, BODY), (r, 9.5, False, INK, BODY)]], space=0.98)
page_num(s, 12)

# ============================== 13. CLOSING ==============================
s = slide(NAVY)
text(s, 0.9, 2.6, 11.5, 1.0, [[("Mechanism over scale.", 40, True, WHITE, HEAD)]])
text(s, 0.95, 3.7, 11.5, 0.6,
     [[("A rigorously-validated associative-memory SSM — every number from a real run.",
        16, False, ICE, BODY)]])
text(s, 0.95, 4.7, 11.5, 0.5,
     [[("github.com/charansaiponnada/Genomic-FM", 14, True, TEAL, BODY)]])
for i, c in enumerate((TEAL, BLUE, AMBER)):
    dot(s, 0.97 + i * 0.34, 4.3, c, 0.18)

out = os.path.join(HERE, "Mnemosyne_Research_Deck.pptx")
prs.save(out)
print("wrote", out, "with", len(prs.slides._sldIdLst), "slides")
